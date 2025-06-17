from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

# Slide 1: Título
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Gestão DOM"
slide.placeholders[1].text = (
    "Transformando a gestão de operações e pessoas\n"
    "Seu nome – Junho/2024"
)

# Slide 2: Oportunidade de Negócio
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Oportunidade de Negócio"
slide.placeholders[1].text = (
    "Vivemos um cenário de processos manuais e descentralizados, que limitam o crescimento e a inovação.\n\n"
    "A Gestão DOM surge para romper barreiras, trazendo agilidade, integração e inteligência para o centro do negócio."
)

# Slide 3: Proposta de Valor
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Proposta de Valor"
slide.placeholders[1].text = (
    "• Centralize informações e tome decisões com confiança.\n"
    "• Automatize rotinas e libere tempo para o que realmente importa.\n"
    "• Transforme dados em resultados.\n\n"
    "Frase de impacto: 'Sua equipe no centro da inovação.'"
)

# Slide 4: Diferenciais Estratégicos
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Diferenciais Estratégicos"
slide.placeholders[1].text = (
    "• Flexibilidade e modularidade: evolua sem limites.\n"
    "• Internacionalização nativa: pronto para crescer.\n"
    "• Experiência do usuário: simples, acessível, encantadora.\n"
    "• Segurança e compliance: confiança para o futuro.\n\n"
    "Frase de impacto: 'Tecnologia que aproxima pessoas e potencializa resultados.'"
)

# Slide 5: Inovação Tecnológica
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Inovação Tecnológica"
slide.placeholders[1].text = (
    "• Stack moderna: Next.js, React, Material UI, Prisma, Postgres.\n"
    "• Código limpo, padronizado e testado.\n"
    "• Integração contínua e governança de dados.\n\n"
    "Frase de impacto: 'Soluções robustas para desafios reais.'"
)

# Slide 6: Impacto no Negócio
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Impacto no Negócio"
slide.placeholders[1].text = (
    "• Redução de custos operacionais\n"
    "• Aumento da produtividade\n"
    "• Decisões mais rápidas e assertivas\n"
    "• Base sólida para inovação contínua\n\n"
    "Frase de impacto: 'Mais do que tecnologia, entregamos valor.'"
)

# Slide 7: Resultados Alcançados
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Resultados Alcançados"
slide.placeholders[1].text = (
    "• Implantação rápida e sem retrabalho\n"
    "• Feedback positivo dos usuários\n"
    "• Pronto para novas integrações e expansão\n\n"
    "Frase de impacto: 'Cada conquista é um passo rumo ao futuro.'"
)

# Slide 8: Telas do Sistema
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Slide em branco para inserir prints
slide.shapes.title.text = "Telas do Sistema"
# Instrução para o usuário
left = top = width = height = Pt(50)
textbox = slide.shapes.add_textbox(left, Pt(100), Pt(800), Pt(100))
tf = textbox.text_frame
tf.text = (
    "Insira aqui os prints das principais telas do sistema.\n"
    "Sugestão: Cadastro de empregado, painel administrativo, dashboard de operações.\n\n"
    "Frase de impacto: 'Visualize o futuro da sua operação em uma única tela.'"
)

# Slide 9: Próximos Passos Estratégicos
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Próximos Passos Estratégicos"
slide.placeholders[1].text = (
    "• Novas integrações (ex: folha de pagamento, BI)\n"
    "• Expansão para outras unidades/países\n"
    "• Evolução contínua baseada em feedback e dados\n\n"
    "Frase de impacto: 'O amanhã começa hoje.'"
)

# Slide 10: Conclusão
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusão"
slide.placeholders[1].text = (
    "Gestão DOM é mais que um sistema: é um diferencial competitivo para o negócio.\n"
    "Pronto para suportar o crescimento e a transformação digital da empresa.\n\n"
    "Frase de impacto: 'Junte-se a nós nessa jornada de transformação.'"
)

# Slide 11: Perguntas e Discussão
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Perguntas e Discussão"
slide.placeholders[1].text = (
    "Dúvidas? Sugestões? Vamos juntos construir o futuro da gestão!"
)

prs.save("apresentacao_gestao_dom_estrategica.pptx") 