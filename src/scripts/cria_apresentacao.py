from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Criação da apresentação
prs = Presentation()

# Slide 1: Título
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Gestão DOM"
slide.placeholders[1].text = "Sistema de Gestão de Operações e Pessoas\nApresentação do Projeto\nSeu Nome - Junho/2024"

# Slide 2: Introdução
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Introdução"
content = (
    "O projeto Gestão DOM foi criado para otimizar e digitalizar processos de gestão de operações e pessoas, "
    "trazendo mais eficiência, controle e integração para a empresa.\n\n"
    "Objetivo: Automatizar rotinas, centralizar informações e facilitar o acesso aos dados."
)
slide.placeholders[1].text = content

# Slide 3: Escopo e Funcionalidades
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Escopo e Funcionalidades"
content = (
    "- Cadastro e gestão de empregados\n"
    "- Controle de operações e tarefas\n"
    "- Fluxo de login e autenticação seguro\n"
    "- Painel administrativo responsivo\n"
    "- Integração com banco de dados Postgres\n"
    "- Internacionalização de mensagens"
)
slide.placeholders[1].text = content

# Slide 4: Arquitetura e Tecnologias
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Arquitetura e Tecnologias"
content = (
    "- Next.js (frontend e backend)\n"
    "- React\n"
    "- Material UI (tema customizado)\n"
    "- Prisma ORM\n"
    "- PostgreSQL\n"
    "- ESLint, TypeScript strict mode\n"
    "- Testes automatizados (Jest, React Testing Library)"
)
slide.placeholders[1].text = content

# Slide 5: Fluxo de Trabalho
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Fluxo de Trabalho"
content = (
    "1. Usuário acessa o sistema via login seguro\n"
    "2. Realiza cadastro ou consulta de empregados\n"
    "3. Gerencia operações e tarefas\n"
    "4. Todas as ações validadas por formulários acessíveis e internacionalizados\n"
    "5. Dados persistidos via Prisma/Postgres"
)
slide.placeholders[1].text = content

# Slide 6: Boas Práticas e Regras
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Boas Práticas e Regras"
content = (
    "- Tipagem estrita (proibido 'any')\n"
    "- Organização modular dos arquivos\n"
    "- Uso de aliases '@/src' nos imports\n"
    "- Mensagens centralizadas e internacionalizáveis\n"
    "- Formulários acessíveis e responsivos\n"
    "- Testes automatizados obrigatórios\n"
    "- Proibido arquivos duplicados ou de backup"
)
slide.placeholders[1].text = content

# Slide 7: Desafios e Soluções
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Desafios e Soluções"
content = (
    "- Padronização de formulários reutilizáveis\n"
    "- Divisão de arquivos grandes por responsabilidade\n"
    "- Integração segura com API e banco\n"
    "- Garantia de acessibilidade e responsividade\n"
    "- Centralização de mensagens para facilitar tradução"
)
slide.placeholders[1].text = content

# Slide 8: Resultados e Benefícios
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Resultados e Benefícios"
content = (
    "- Código limpo, padronizado e fácil de manter\n"
    "- Experiência do usuário aprimorada\n"
    "- Facilidade para evoluir e integrar novas funções\n"
    "- Redução de erros e retrabalho\n"
    "- Base pronta para crescimento do sistema"
)
slide.placeholders[1].text = content

# Slide 9: Próximos Passos
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Próximos Passos"
content = (
    "- Novas integrações (ex: folha de pagamento)\n"
    "- Mais testes automatizados\n"
    "- Melhorias de UX e acessibilidade\n"
    "- Otimização de performance\n"
    "- Feedback contínuo dos usuários"
)
slide.placeholders[1].text = content

# Slide 10: Demonstração (opcional)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Demonstração"
content = (
    "(Inserir prints ou vídeo do sistema em funcionamento)\n"
    "Exemplo: Cadastro de empregado, login, painel administrativo."
)
slide.placeholders[1].text = content

# Slide 11: Perguntas
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Perguntas"
content = "Dúvidas? Sugestões? Estou à disposição para responder!"
slide.placeholders[1].text = content

# Salvar apresentação
prs.save("apresentacao_gestao_dom.pptx") 